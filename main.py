import random
import string
from fastapi import FastAPI, File, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from io import BytesIO
import fitz
import pytesseract
from PIL import Image
import spacy
from docx import Document
from pptx import Presentation

# Load the NLP model for sentence parsing
nlp = spacy.load("en_core_web_sm")

# Initialize FastAPI app
app = FastAPI()

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Change to your frontend's origin in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Function to extract text from a PDF file and images (OCR)
def extract_text_from_pdf_and_images(pdf_file: UploadFile):
    pdf_bytes = pdf_file.file.read()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
        for img in page.get_images(full=True):  # OCR for images
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image = Image.open(BytesIO(image_bytes))
            text += pytesseract.image_to_string(image)
    return text

# Function to extract text from a Word document
def extract_text_from_word(word_file: UploadFile):
    word_bytes = word_file.file.read()
    doc = Document(BytesIO(word_bytes))
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

# Function to extract text from a PowerPoint file
def extract_text_from_ppt(ppt_file: UploadFile):
    ppt_bytes = ppt_file.file.read()
    prs = Presentation(BytesIO(ppt_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to extract text from various file types
def extract_text_from_file(file: UploadFile):
    file_extension = file.filename.split('.')[-1].lower()
    if file_extension == "pdf":
        return extract_text_from_pdf_and_images(file)
    elif file_extension == "docx":
        return extract_text_from_word(file)
    elif file_extension == "pptx":
        return extract_text_from_ppt(file)
    else:
        raise ValueError("Unsupported file type. Please upload a PDF, Word, or PowerPoint file.")

# Function to identify sentences from extracted text
def identify_key_sentences(text):
    doc = nlp(text)
    return [sent.text for sent in doc.sents]

# Function to generate random multiple-choice questions
def generate_mcq(question, correct_answer, all_answers):
    choices = [correct_answer] + random.sample(all_answers, 3)  # Add 3 incorrect answers
    random.shuffle(choices)  # Shuffle the order
    return f"{question}\n" + "\n".join(f"{chr(65+i)}) {choice}" for i, choice in enumerate(choices))

# Function to generate short-answer questions
def generate_short_answer(question):
    return f"{question}\n(Provide your answer here.)"

# Function to generate true/false questions
def generate_true_false(question, correct_answer):
    return f"{question}\n- True / False\nCorrect Answer: {correct_answer}"

# Function to analyze text and generate questions based on identified entities
def generate_questions_from_text(text):
    questions = []
    sentences = identify_key_sentences(text)

    for sentence in sentences:
        if random.random() < 0.33:  # MCQ
            question_text = f"What does the following mean: '{sentence}'?"
            questions.append(generate_mcq(question_text, sentence, ["Example A", "Example B", "Example C"]))
        elif random.random() < 0.5:  # Short Answer
            questions.append(generate_short_answer(f"Explain: '{sentence}'."))
        else:  # True/False
            questions.append(generate_true_false(f"Is the following correct? '{sentence}'", "True" if random.random() > 0.5 else "False"))
    return questions

@app.post("/generate_exam/")
async def generate_exam_from_file(file: UploadFile = File(...)):
    try:
        extracted_text = extract_text_from_file(file)
        questions = generate_questions_from_text(extracted_text)
        return {"questions": questions}
    except ValueError as e:
        return {"error": str(e)}
    except Exception as e:
        print(f"Unexpected Error: {e}")
        return {"error": "An unexpected error occurred."}
